import argparse
import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN

#CUSTOMIZE HERE THE LOOK OF THE PPTX
FONT = "Arial"
FONT_SIZE = Inches(0.55) # = 55 in ppt
######################################



# Set up the command-line argument parser
parser = argparse.ArgumentParser(description='Scrape text from a web page and create a PowerPoint presentation.')
parser.add_argument('url', type=str, help='the URL of the web page to scrape')
# Parse the command-line arguments
args = parser.parse_args()

# Get the HTML content of the page
response = requests.get(args.url)
html_content = response.content

# Parse the HTML content with Beautiful Soup
soup = BeautifulSoup(html_content, 'html.parser')

# Find all the elements with the class "resized-text" and extract their HTML content
element = str(soup.find(class_='resized-text'))

# Split the HTML content into groups separated by 2 <br> tags
text_list = []

text_parts = element.split('<br/> <br/>')
if len(text_parts) == 1:
    text_parts = element.split('<br/>\n<br/>')

# Extract the text content of each group of HTML tags
text_parts = [BeautifulSoup(part, 'html.parser').get_text().strip() for part in text_parts]
# Remove any empty groups of text
text_parts = [part for part in text_parts if part]
# Add the non-empty groups of text to the list
text_list.extend(text_parts)

# Create a new PowerPoint presentation and add each group of text to a new slide
prs = Presentation()
for text in text_list:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)
    text_frame = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(6))
    text_frame.text = text
    text_frame.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE # center text vertically
    text_frame.text_frame.word_wrap = True # wrap text within the textbox
    text_frame.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE # auto-size the text to fit the shape
    text_frame.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER # center text horizontally

    for paragraph in text_frame.text_frame.paragraphs:
        paragraph.font.color.rgb = RGBColor(255, 255, 255) # set text color to white
        paragraph.font.name = 'Arial' # set font to Times New Roman
        paragraph.font.size = Inches(0.55) # set font size to 55 points
        paragraph.font.bold = True # set font to bold
# Save the PowerPoint presentation
prs.save("output.pptx")
